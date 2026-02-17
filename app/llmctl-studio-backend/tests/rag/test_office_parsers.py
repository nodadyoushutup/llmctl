import os
import sys
import tempfile
import unittest
import importlib.util
from pathlib import Path

REPO_ROOT = Path(__file__).resolve().parents[4]
STUDIO_SRC = REPO_ROOT / "app" / "llmctl-studio" / "src"
STUDIO_APP_ROOT = REPO_ROOT / "app" / "llmctl-studio"
if str(STUDIO_SRC) not in sys.path:
    sys.path.insert(0, str(STUDIO_SRC))
if str(STUDIO_APP_ROOT) not in sys.path:
    sys.path.insert(0, str(STUDIO_APP_ROOT))

from rag.engine.office_parsers import parse_docx, parse_pptx, parse_xlsx

_HELPERS_SPEC = importlib.util.spec_from_file_location(
    "rag_test_helpers",
    STUDIO_APP_ROOT / "tests" / "rag" / "helpers.py",
)
if _HELPERS_SPEC is None or _HELPERS_SPEC.loader is None:  # pragma: no cover
    raise RuntimeError("Failed to load rag test helpers.")
_HELPERS_MODULE = importlib.util.module_from_spec(_HELPERS_SPEC)
_HELPERS_SPEC.loader.exec_module(_HELPERS_MODULE)
test_config = _HELPERS_MODULE.test_config

try:
    import docx
except ImportError:
    docx = None

try:
    import pptx
except ImportError:
    pptx = None

try:
    import openpyxl
except ImportError:
    openpyxl = None


class OfficeParserTests(unittest.TestCase):
    @unittest.skipUnless(docx, "python-docx is required")
    def test_docx_parser(self):
        with tempfile.TemporaryDirectory() as tmpdir:
            path = Path(tmpdir) / "sample.docx"
            document = docx.Document()
            document.add_heading("Title", level=1)
            document.add_paragraph("Hello world")
            document.save(path)
            parsed = parse_docx(path, test_config(Path(tmpdir)))
            self.assertIsNotNone(parsed)
            self.assertEqual(parsed.doc_type, "docx")

    @unittest.skipUnless(pptx, "python-pptx is required")
    def test_pptx_parser(self):
        with tempfile.TemporaryDirectory() as tmpdir:
            path = Path(tmpdir) / "sample.pptx"
            pres = pptx.Presentation()
            slide = pres.slides.add_slide(pres.slide_layouts[5])
            title = slide.shapes.title
            if title is not None:
                title.text = "Hello"
            pres.save(path)
            parsed = parse_pptx(path, test_config(Path(tmpdir)))
            self.assertIsNotNone(parsed)
            self.assertEqual(parsed.doc_type, "pptx")

    @unittest.skipUnless(openpyxl, "openpyxl is required")
    def test_xlsx_parser(self):
        with tempfile.TemporaryDirectory() as tmpdir:
            path = Path(tmpdir) / "sample.xlsx"
            wb = openpyxl.Workbook()
            ws = wb.active
            ws["A1"] = "Header"
            ws["A2"] = "Row1"
            wb.save(path)
            parsed = parse_xlsx(path, test_config(Path(tmpdir)))
            self.assertIsNotNone(parsed)
            self.assertEqual(parsed.doc_type, "xlsx")


if __name__ == "__main__":
    unittest.main()
