from __future__ import annotations

import hashlib
from pathlib import Path

from rag.engine.config import RagConfig, max_file_bytes_for
from rag.engine.pipeline import ParsedDocument

try:
    import docx
except ImportError:  # pragma: no cover
    docx = None

try:
    import pptx
except ImportError:  # pragma: no cover
    pptx = None

try:
    import openpyxl
except ImportError:  # pragma: no cover
    openpyxl = None


def _require_docx():
    if docx is None:
        raise RuntimeError("python-docx is required to parse DOCX")


def _require_pptx():
    if pptx is None:
        raise RuntimeError("python-pptx is required to parse PPTX")


def _require_xlsx():
    if openpyxl is None:
        raise RuntimeError("openpyxl is required to parse XLSX")


def _base_source(path: Path):
    stat = path.stat()
    data = path.read_bytes()
    return {
        "path": str(path),
        "size": stat.st_size,
        "mtime": stat.st_mtime,
        "file_hash": hashlib.sha1(data).hexdigest(),
    }


def parse_docx(path: Path, config: RagConfig) -> ParsedDocument | None:
    _require_docx()
    if path.stat().st_size > max_file_bytes_for(config, "docx"):
        return None
    doc = docx.Document(path)
    spans = []
    current_heading = None
    buffer = []
    start_line = 1
    line_no = 1

    for paragraph in doc.paragraphs:
        text = paragraph.text
        style = paragraph.style.name if paragraph.style else ""
        is_heading = style.lower().startswith("heading")
        if is_heading:
            if buffer:
                spans.append(
                    {
                        "text": "\n".join(buffer),
                        "start_line": start_line,
                        "end_line": line_no - 1,
                        "metadata": {
                            "section": current_heading or "<untitled>",
                        },
                    }
                )
                buffer = []
            current_heading = text.strip() or "<untitled>"
            start_line = line_no
        if text.strip():
            buffer.append(text)
            line_no += 1

    if buffer:
        spans.append(
            {
                "text": "\n".join(buffer),
                "start_line": start_line,
                "end_line": line_no - 1,
                "metadata": {
                    "section": current_heading or "<untitled>",
                },
            }
        )

    content = "\n".join([span["text"] for span in spans])
    return ParsedDocument(
        content=content,
        doc_type="docx",
        language=None,
        source=_base_source(path),
        structural_hints={"spans": spans} if spans else {},
    )


def parse_pptx(path: Path, config: RagConfig) -> ParsedDocument | None:
    _require_pptx()
    if path.stat().st_size > max_file_bytes_for(config, "pptx"):
        return None
    pres = pptx.Presentation(path)
    spans = []
    parts = []
    for idx, slide in enumerate(pres.slides, start=1):
        text_runs = []
        for shape in slide.shapes:
            if not hasattr(shape, "text"):
                continue
            text = shape.text
            if text.strip():
                text_runs.append(text.strip())
        slide_text = "\n".join(text_runs)
        if not slide_text:
            continue
        parts.append(slide_text)
        spans.append(
            {
                "text": slide_text,
                "start_line": None,
                "end_line": None,
                "metadata": {
                    "slide": idx,
                },
            }
        )

    content = "\n\n".join(parts)
    return ParsedDocument(
        content=content,
        doc_type="pptx",
        language=None,
        source=_base_source(path),
        structural_hints={"spans": spans} if spans else {},
    )


def parse_xlsx(path: Path, config: RagConfig) -> ParsedDocument | None:
    _require_xlsx()
    if path.stat().st_size > max_file_bytes_for(config, "xlsx"):
        return None
    wb = openpyxl.load_workbook(path, data_only=True, read_only=True)
    spans = []
    parts = []
    for sheet in wb.worksheets:
        rows = []
        for row in sheet.iter_rows():
            values = []
            for cell in row:
                if cell.value is None:
                    values.append("")
                else:
                    values.append(str(cell.value))
            if any(values):
                rows.append("\t".join(values))
        sheet_text = "\n".join(rows)
        if not sheet_text.strip():
            continue
        parts.append(sheet_text)
        spans.append(
            {
                "text": sheet_text,
                "start_line": None,
                "end_line": None,
                "metadata": {
                    "sheet": sheet.title,
                },
            }
        )

    content = "\n\n".join(parts)
    return ParsedDocument(
        content=content,
        doc_type="xlsx",
        language=None,
        source=_base_source(path),
        structural_hints={"spans": spans} if spans else {},
    )
