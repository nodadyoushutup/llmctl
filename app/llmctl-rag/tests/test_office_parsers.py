import os
import sys
import tempfile
import unittest
from pathlib import Path

ROOT = Path(__file__).resolve().parents[1]
sys.path.insert(0, str(ROOT))

from office_parsers import parse_docx, parse_pptx, parse_xlsx
from tests.helpers import test_config

try:
    import docx
except ImportError:
    docx = None

try:
    import pptx
except ImportError:
    pptx = None

try:
    import openpyxl
except ImportError:
    openpyxl = None


class OfficeParserTests(unittest.TestCase):
    @unittest.skipUnless(docx, "python-docx is required")
    def test_docx_parser(self):
        with tempfile.TemporaryDirectory() as tmpdir:
            path = Path(tmpdir) / "sample.docx"
            document = docx.Document()
            document.add_heading("Title", level=1)
            document.add_paragraph("Hello world")
            document.save(path)
            parsed = parse_docx(path, test_config(Path(tmpdir)))
            self.assertIsNotNone(parsed)
            self.assertEqual(parsed.doc_type, "docx")

    @unittest.skipUnless(pptx, "python-pptx is required")
    def test_pptx_parser(self):
        with tempfile.TemporaryDirectory() as tmpdir:
            path = Path(tmpdir) / "sample.pptx"
            pres = pptx.Presentation()
            slide = pres.slides.add_slide(pres.slide_layouts[5])
            title = slide.shapes.title
            if title is not None:
                title.text = "Hello"
            pres.save(path)
            parsed = parse_pptx(path, test_config(Path(tmpdir)))
            self.assertIsNotNone(parsed)
            self.assertEqual(parsed.doc_type, "pptx")

    @unittest.skipUnless(openpyxl, "openpyxl is required")
    def test_xlsx_parser(self):
        with tempfile.TemporaryDirectory() as tmpdir:
            path = Path(tmpdir) / "sample.xlsx"
            wb = openpyxl.Workbook()
            ws = wb.active
            ws["A1"] = "Header"
            ws["A2"] = "Row1"
            wb.save(path)
            parsed = parse_xlsx(path, test_config(Path(tmpdir)))
            self.assertIsNotNone(parsed)
            self.assertEqual(parsed.doc_type, "xlsx")


if __name__ == "__main__":
    unittest.main()
